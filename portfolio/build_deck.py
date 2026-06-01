"""KBO 베이지안 포트폴리오 발표 덱 생성 (python-pptx)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
CH = os.path.join(HERE, "charts")

NAVY = RGBColor(0x1E, 0x27, 0x61)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
GOLD = RGBColor(0xE8, 0xB6, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x22, 0x26, 0x2C)
MUTE = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF4, 0xF5, 0xF7)
FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = bg
    return s


def _setfont(run, name=FONT):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.1):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs, str):
        runs = [(runs, 18, INK, False)]
    for i, (t, sz, col, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.color.rgb = col
        r.font.bold = bold
        _setfont(r)
    return tb


def rect(s, x, y, w, h, fill, line=None, round_=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    return shp


def card(s, x, y, w, h, lines, fill=WHITE, border=ICE):
    rect(s, x, y, w, h, fill, line=border, round_=True)
    text(s, x + 0.25, y + 0.18, w - 0.5, h - 0.36, lines, anchor=MSO_ANCHOR.MIDDLE)


def pic(s, name, x, y, w):
    p = s.shapes.add_picture(os.path.join(CH, name), Inches(x), Inches(y), width=Inches(w))
    return p


def pic_centered(s, name, top, height):
    from PIL import Image
    iw, ih = Image.open(os.path.join(CH, name)).size
    w = height * iw / ih
    p = s.shapes.add_picture(os.path.join(CH, name), Inches((SW - w) / 2), Inches(top),
                             height=Inches(height))
    return p


def stat_row(s, items, y, h=1.05):
    n = len(items)
    gap = 0.3
    total = SW - 1.4
    w = (total - gap * (n - 1)) / n
    for i, (num, lab, col) in enumerate(items):
        x = 0.7 + i * (w + gap)
        rect(s, x, y, w, h, LIGHT, line=ICE, round_=True)
        text(s, x + 0.2, y + 0.12, w - 0.4, h - 0.24,
             [(num, 30, col, True), (lab, 12, MUTE, False)],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line=1.0)


def caption(s, t, y=7.0):
    text(s, 0.7, y, SW - 1.4, 0.4, [(t, 11, MUTE, False)])


# ───────────────────────── S1 Title ─────────────────────────
s = slide(NAVY)
rect(s, 0, 0, SW, 0.18, GOLD)
text(s, 0.9, 2.2, 11.5, 1.6,
     [("KBO 베이지안 추론 파이프라인", 46, WHITE, True)])
text(s, 0.9, 3.7, 11.5, 1.4,
     [("시즌 초 타율의 평균회귀를 베이지안으로 보정하고,", 20, ICE, False),
      ("매일 자동으로 수집·갱신하는 야구 통계 추론 시스템", 20, ICE, False)], line=1.3)
text(s, 0.9, 6.2, 11.5, 0.8,
     [("데이터 분석 포트폴리오   ·   Live  kbo-bayes.vercel.app   ·   GitHub  github.com/icedo724/kbo-bayes",
       13, ICE, False)])

# ───────────────────────── S2 Problem ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.5, 12, 1.0, [("시즌 초 타율은 믿을 수 없다", 36, NAVY, True)])
text(s, 0.7, 1.6, 5.5, 4.5,
     [("20타수에서 8안타를 친 선수의 타율은 .400.", 18, INK, False),
      ("하지만 이 선수가 진짜 4할 타자일 가능성은 낮다 —", 18, INK, False),
      ("대부분 시간이 지나면 평균으로 회귀한다.", 18, INK, False),
      ("", 8, INK, False),
      ("'관측 타율을 그대로 쓰는 것'이 베이스라인인데,", 18, INK, False),
      ("표본이 적을수록 과대·과소 추정이 심하다.", 18, INK, False),
      ("", 8, INK, False),
      ("→ 표본 크기를 반영해 '실력'을 추정해야 한다.", 18, GOLD, True)], line=1.35)
card(s, 6.6, 1.8, 3.0, 1.7,
     [("관측 타율", 13, MUTE, False), (".400", 40, RGBColor(0xCF, 0x24, 0x2E), True),
      ("20타수 8안타", 12, MUTE, False)], fill=LIGHT)
card(s, 9.9, 1.8, 3.0, 1.7,
     [("리그 평균", 13, MUTE, False), (".254", 40, NAVY, True),
      ("진짜 실력은 이 사이 어딘가", 12, MUTE, False)], fill=LIGHT)
card(s, 6.6, 3.9, 6.3, 2.0,
     [("핵심 질문", 15, GOLD, True),
      ("\"표본이 적은 선수의 타율을 어떻게 보정해야", 17, INK, False),
      ("시즌 잔여 성적을 더 잘 예측할 수 있는가?\"", 17, INK, False)], fill=WHITE)

# ───────────────────────── S3 Approach ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.5, 12, 1.0, [("접근: Beta-Binomial 베이지안 shrinkage", 34, NAVY, True)])
card(s, 0.7, 1.7, 11.9, 1.5,
     [("H ~ Binomial(AB, θ)     θ ~ Beta(α, β)     ⇒     θ | data ~ Beta(α+H, β+AB−H)",
       22, NAVY, True),
      ("켤레 사전분포 → 닫힌형 해. MCMC 없이 매일 가볍게 갱신", 13, MUTE, False)],
     fill=LIGHT, border=ICE)
steps = [
    ("①  리그 prior", "전체 타자 분포에서\nempirical Bayes(MoM)로\nBeta(α,β) 추정", NAVY),
    ("②  사후 갱신", "선수의 H/AB를 더해\n사후분포로 이동.\n표본 적으면 prior에 가까움", GOLD),
    ("③  수축 효과", "타석이 쌓일수록\n관측값으로 수렴,\n신뢰구간도 좁아짐", NAVY),
]
for i, (h, b, col) in enumerate(steps):
    x = 0.7 + i * 4.07
    rect(s, x, 3.6, 3.8, 2.8, WHITE, line=ICE, round_=True)
    text(s, x + 0.3, 3.85, 3.2, 0.6, [(h, 19, col, True)])
    text(s, x + 0.3, 4.6, 3.2, 1.7, [(b, 15, INK, False)], line=1.3)

# ───────────────────────── S4 Result: shrinkage ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.45, 12, 0.9, [("결과 ① 표본이 적을수록 강하게 보정된다", 32, NAVY, True)])
stat_row(s, [
    (".400→.262", "5타수 신인 (−.138)", RGBColor(0xCF, 0x24, 0x2E)),
    (".000→.248", "2타수 무안타 (+.248)", NAVY),
    (".274→.268", "223타수 주전 (−.006)", RGBColor(0x1A, 0x7F, 0x37)),
], y=1.35)
pic_centered(s, "shrinkage.png", top=2.55, height=4.35)

# ───────────────────────── S5 Result: trajectory ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.45, 12, 0.9, [("결과 ② 시즌이 흐르며 추정이 수렴한다", 32, NAVY, True)])
text(s, 0.7, 1.3, 12, 0.6,
     [("초반 관측치(회색)는 크게 출렁이지만 베이지안 추정(네이비)은 안정적이며, 타석이 쌓일수록 신뢰구간이 좁아진다.",
       14, MUTE, False)])
pic_centered(s, "trajectory.png", top=2.0, height=4.9)

# ───────────────────────── S6 Validation ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.45, 12, 0.9, [("정직한 검증: 다년 walk-forward", 32, NAVY, True)])
text(s, 0.7, 1.3, 6.0, 5.0,
     [("평가 원칙", 16, GOLD, True),
      ("• accuracy 금지 → log loss / Brier", 16, INK, False),
      ("• 항상 베이스라인(관측 타율)과 비교", 16, INK, False),
      ("• 시즌별 명단·시점별 prior 재추정으로", 16, INK, False),
      ("   look-ahead·생존편향 차단", 16, INK, False),
      ("", 8, INK, False),
      ("발견 (정직하게 보고)", 16, GOLD, True),
      ("MLE prior는 시즌 초 소표본에서 발산 →", 15, INK, False),
      ("모두 평균으로 과수축. 그래서 안정적인", 15, INK, False),
      ("MoM을 채택했다.", 15, INK, False)], line=1.3)
pic(s, "multiseason.png", 6.9, 1.5, 6.0)
text(s, 6.9, 6.55, 6.0, 0.5,
     [("5개 시즌 × 6시점 = 30/30 시점에서 베이지안 우월", 13, NAVY, True)], align=PP_ALIGN.CENTER)

# ───────────────────────── S7 Playoff ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.45, 12, 0.9, [("확장: 가을야구 진출 확률", 32, NAVY, True)])
text(s, 0.7, 1.4, 5.6, 5.0,
     [("같은 베이지안 아이디어를 팀에 적용", 17, GOLD, True),
      ("• 팀 승률도 Beta-Binomial로 .500 보정", 16, INK, False),
      ("• 잔여 경기를 2만 회 몬테카를로", 16, INK, False),
      ("• 5위 이내 비율 = 진출 확률", 16, INK, False),
      ("", 8, INK, False),
      ("검증 (2025)", 17, GOLD, True),
      ("불확실성이 큰 시즌 초·중반에 베이스라인", 15, INK, False),
      ("(현 승률 유지)보다 Brier 우월 (4/5 시점).", 15, INK, False),
      ("막바지엔 순위 확정으로 결정론이 정확 —", 15, INK, False),
      ("이 한계도 그대로 보고한다.", 15, INK, False)], line=1.3)
pic(s, "playoff.png", 6.5, 1.4, 6.4)

# ───────────────────────── S8 Architecture ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.5, 12, 0.9, [("시스템 아키텍처: 매일 자동으로 도는 추론", 32, NAVY, True)])
flow = [
    ("수집", "KBO 공식\nrequests + read_html\nrobots.txt 준수"),
    ("저장", "Supabase(Postgres)\nservice=쓰기 / anon=읽기\nRLS 정책"),
    ("추정", "동결 prior로\n사후분포 갱신\n(online 읽기 전용)"),
    ("시각화", "Next.js + Vercel\n포털형 대시보드\n(계산 없이 SELECT)"),
]
bw = 2.85
for i, (h, b) in enumerate(flow):
    x = 0.7 + i * (bw + 0.32)
    rect(s, x, 2.0, bw, 1.9, NAVY if i % 2 == 0 else RGBColor(0x2C, 0x36, 0x70), round_=True)
    text(s, x + 0.2, 2.2, bw - 0.4, 0.6, [(h, 20, GOLD, True)], align=PP_ALIGN.CENTER)
    text(s, x + 0.2, 2.85, bw - 0.4, 1.0, [(b, 13, WHITE, False)],
         align=PP_ALIGN.CENTER, line=1.2)
    if i < 3:
        text(s, x + bw - 0.02, 2.55, 0.36, 0.6, [("→", 24, MUTE, True)], align=PP_ALIGN.CENTER)
rect(s, 0.7, 4.5, SW - 1.4, 1.9, LIGHT, line=ICE, round_=True)
text(s, 1.0, 4.7, SW - 2.0, 1.6,
     [("설계 원칙", 16, GOLD, True),
      ("• 오프라인(검증·동결) / 온라인(매일 예측) 물리적 분리 — 자동화가 모델을 건드리지 않음", 15, INK, False),
      ("• look-ahead bias 차단 — 모든 추정은 그 시점까지의 누적 데이터로만", 15, INK, False),
      ("• GitHub Actions cron(매일 새벽) — 수집→갱신→저장 완전 자동", 15, INK, False)], line=1.35)

# ───────────────────────── S9 Dashboard ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.5, 12, 0.9, [("라이브 대시보드 (포털형)", 32, NAVY, True)])
feats = [
    ("홈", "팀 카드 · 최근 경기"),
    ("팀 전력 분석", "팀 타자 shrinkage · 지표"),
    ("선수 상세", "시즌 궤적 · 신뢰구간"),
    ("리그 분석", "필터 · 선수 비교"),
    ("일정·결과", "매치업·승패 재구성"),
    ("진출 확률", "팀별 가을야구 %"),
]
for i, (h, b) in enumerate(feats):
    r, c = divmod(i, 3)
    x = 0.7 + c * 4.07
    y = 1.7 + r * 1.55
    rect(s, x, y, 3.8, 1.35, WHITE, line=ICE, round_=True)
    text(s, x + 0.25, y + 0.2, 3.3, 0.5, [(h, 18, NAVY, True)])
    text(s, x + 0.25, y + 0.78, 3.3, 0.45, [(b, 13, MUTE, False)])
rect(s, 0.7, 5.1, SW - 1.4, 1.4, NAVY, round_=True)
text(s, 0.7, 5.45, SW - 1.4, 0.8,
     [("kbo-bayes.vercel.app", 28, WHITE, True)], align=PP_ALIGN.CENTER)

# ───────────────────────── S10 Limits & Lessons ─────────────────────────
s = slide(WHITE)
text(s, 0.7, 0.5, 12, 0.9, [("한계 & 배운 점", 34, NAVY, True)])
rect(s, 0.7, 1.7, 5.9, 4.6, LIGHT, line=ICE, round_=True)
text(s, 1.0, 1.95, 5.3, 4.2,
     [("한계 (정직하게)", 18, GOLD, True),
      ("• 타자 데이터만 — 투수/팀 종합은 향후", 15, INK, False),
      ("• 일정/스코어: KBO API가 robots(/ws/)로", 15, INK, False),
      ("   차단 → 합법 데이터로 매치업·승패만 재구성", 15, INK, False),
      ("• 진출확률은 팀 독립 시뮬(상대전적 미반영)", 15, INK, False),
      ("• 타율은 정보량 한계 → OBP/wOBA 확장 예정", 15, INK, False)], line=1.35)
rect(s, 6.8, 1.7, 5.85, 4.6, NAVY, round_=True)
text(s, 7.1, 1.95, 5.25, 4.2,
     [("배운 점", 18, GOLD, True),
      ("• 데이터 누수 방지는 코드 구조로 강제해야", 15, WHITE, False),
      ("   안전하다 (오프라인/온라인 분리)", 15, WHITE, False),
      ("• '못 이기면 못 이긴다고 보고'하는 정직한", 15, WHITE, False),
      ("   검증이 신뢰를 만든다", 15, WHITE, False),
      ("• robots.txt 준수 — 되는 것과 해도 되는 것은", 15, WHITE, False),
      ("   다르다", 15, WHITE, False),
      ("• 분석을 '운영되는 시스템'으로 만드는 경험", 15, WHITE, False)], line=1.35)

# ───────────────────────── S11 Closing ─────────────────────────
s = slide(NAVY)
rect(s, 0, 0, SW, 0.18, GOLD)
text(s, 0.9, 1.6, 11.5, 1.2, [("매일 스스로 추론하는 야구 분석 시스템", 38, WHITE, True)])
text(s, 0.9, 3.0, 11.5, 0.6,
     [("Python · NumPy/SciPy · Supabase(Postgres) · GitHub Actions · Next.js/Vercel", 17, ICE, False)])
card(s, 0.9, 3.9, 5.6, 1.5,
     [("Live Demo", 14, GOLD, True), ("kbo-bayes.vercel.app", 20, WHITE, True)],
     fill=RGBColor(0x2C, 0x36, 0x70), border=NAVY)
card(s, 6.8, 3.9, 5.6, 1.5,
     [("Source", 14, GOLD, True), ("github.com/icedo724/kbo-bayes", 18, WHITE, True)],
     fill=RGBColor(0x2C, 0x36, 0x70), border=NAVY)
text(s, 0.9, 5.9, 11.5, 0.6, [("수집 · 베이지안 추정 · 자동화 · 시각화 — end-to-end", 15, ICE, False)])

OUT = os.path.join(HERE, "KBO_베이지안_포트폴리오.pptx")
prs.save(OUT)
print("saved →", OUT, f"({len(prs.slides._sldIdLst)} slides)")
